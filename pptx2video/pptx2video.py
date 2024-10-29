#!/usr/bin/env python3
"""
.. module:: scripts.pptx2video
   :synopsis: Convert PowerPoint presentations to video with AI-generated voiceovers.

This script converts a PowerPoint presentation (.pptx) into a video (.mp4) by:
- Extracting slides as images (via PDF intermediate).
- Generating voiceover audio using OpenAI TTS API based on slide notes.
- Combining images and audio into video segments per slide.
- Concatenating the slide videos into a final video.

The process is idempotent and utilizes both asynchronous and parallel processing to improve performance.

Usage:
    python pptx2video.py <presentation.pptx> [--force] [--changed]

Requirements:
    - Python 3.7+
    - Install required packages:
        pip install python-pptx Pillow openai pdf2image aiohttp tenacity
    - LibreOffice must be installed and accessible in the system PATH.
    - FFmpeg must be installed and accessible in the system PATH.
    - Set the 'OPENAI_API_KEY' environment variable with your OpenAI API key.

Note:
    Ensure you have sufficient permissions and API quota for OpenAI TTS API.
"""

import argparse
import asyncio
import hashlib
import json
import logging
import multiprocessing
import os
import subprocess
import sys
from concurrent.futures import ProcessPoolExecutor
from typing import List, Optional, Tuple

import aiohttp
import openai
import pdf2image
from PIL import Image
from pptx import Presentation
from tenacity import (
    retry,
    retry_if_exception_type,
    stop_after_attempt,
    wait_exponential,
)

logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

MAX_CONCURRENT_CALLS = 5  # Maximum number of concurrent OpenAI API calls


@retry(
    stop=stop_after_attempt(3),
    wait=wait_exponential(multiplier=1, min=4, max=10),
    retry=retry_if_exception_type((aiohttp.ClientError, openai.OpenAIError)),
)
async def tts_async(input_text: str, model: str, voice: str, api_key: str) -> bytes:
    """Asynchronously convert input text to speech using OpenAI's Text-to-Speech API.

    :param input_text: The text to be converted to speech.
    :param model: The model to use for synthesis (e.g., 'tts-1', 'tts-1-hd').
    :param voice: The voice profile to use (e.g., 'alloy', 'echo', 'fable', etc.).
    :param api_key: OpenAI API key.
    :return: Audio content as bytes.
    :raises ValueError: If API key or input text are empty.
    :raises RuntimeError: If the API call fails.
    """
    if not api_key.strip() or not input_text.strip():
        raise ValueError("API key and input text are required.")

    logger.debug(f"Sending TTS request for text: {input_text[:50]}...")

    async with aiohttp.ClientSession() as session:
        async with session.post(
            "https://api.openai.com/v1/audio/speech",
            headers={"Authorization": f"Bearer {api_key}"},
            json={"model": model.lower(), "voice": voice.lower(), "input": input_text},
        ) as response:
            if response.status != 200:
                error_text = await response.text()
                raise RuntimeError(
                    f"API call failed with status {response.status}: {error_text}"
                )

            audio_content = await response.read()
            logger.debug(f"Received audio content, size: {len(audio_content)} bytes")
            return audio_content


def convert_slide_to_image(args) -> Tuple[str, Optional[str], int]:
    """Convert a single slide from the PDF to an image.

    :param args: Tuple containing pdf_filename, output_dir, slide index i, force_flag.
    :return: Tuple of (image_path, image_hash, slide index)
    """
    pdf_filename, output_dir, i, force_flag = args
    image_path = os.path.join(output_dir, f"slide_{i}.png")

    # Check if image already exists and force_flag is not set
    if not force_flag and os.path.exists(image_path):
        with open(image_path, "rb") as f:
            image_hash = hashlib.sha256(f.read()).hexdigest()
        logger.info(f"Image for slide {i+1} exists, skipping conversion")
        return image_path, image_hash, i

    try:
        images = pdf2image.convert_from_path(
            pdf_filename, first_page=i + 1, last_page=i + 1, dpi=300
        )
        if images:
            images[0].save(image_path, "PNG")
            logger.info(f"Image for slide {i+1} generated")
            # Compute image hash
            with open(image_path, "rb") as f:
                image_hash = hashlib.sha256(f.read()).hexdigest()
        else:
            logger.warning(f"Failed to convert slide {i+1} to image")
            image_hash = None
    except Exception as e:
        logger.error(f"Error converting slide {i+1} to image: {e}")
        image_hash = None

    return image_path, image_hash, i


async def generate_audio_for_slide(
    text: str,
    text_hash: str,
    output_dir: str,
    i: int,
    api_key: str,
    state: dict,
    force_flag: bool,
) -> Optional[str]:
    """Generate audio for a single slide using the TTS function.

    :param text: Text to convert to speech.
    :param text_hash: Hash of the text.
    :param output_dir: Directory to store output files.
    :param i: Slide index.
    :param api_key: OpenAI API key.
    :param state: State dictionary to track progress.
    :param force_flag: If True, forces regeneration of the audio.
    :return: Path to the generated audio file, or None if generation fails.
    """
    slide_audio_filename = os.path.join(output_dir, f"voice_{i}.mp3")

    # Check if audio is up to date
    if (
        not force_flag
        and state["slide_hashes"].get(str(i)) == text_hash
        and os.path.exists(slide_audio_filename)
    ):
        logger.info(f"Audio for slide {i+1} is up to date, skipping generation")
        return slide_audio_filename

    if not text.strip():
        logger.info(f"No notes for slide {i+1}, skipping audio generation")
        return None

    try:
        logger.info(f"Generating audio for slide {i+1}")
        audio_content = await tts_async(text, "tts-1-hd", "echo", api_key)
        with open(slide_audio_filename, "wb") as f:
            f.write(audio_content)
        # Update state in main process
        state["slide_hashes"][str(i)] = text_hash
        return slide_audio_filename
    except Exception as e:
        logger.error(f"Error generating audio for slide {i+1}: {e}")
        return None


def create_video_for_slide(args) -> Optional[str]:
    """Create a video for a single slide by combining image and audio.

    :param args: Tuple containing image_file, audio_file, output_dir, slide index i, slide_content_hash, force_flag.
    :return: Path to the generated video file.
    """
    image_file, audio_file, output_dir, i, slide_content_hash, force_flag = args
    slide_video_filename = os.path.join(output_dir, f"video_{i}.mp4")

    # Check if video already exists and is up to date
    if not force_flag and os.path.exists(slide_video_filename):
        logger.info(f"Video for slide {i+1} exists, skipping creation")
        return slide_video_filename

    logger.info(f"Creating video for slide {i+1}")

    # Get image dimensions and ensure width is even for video encoding
    with Image.open(image_file) as img:
        width, height = img.size
    adjusted_width = width if width % 2 == 0 else width - 1

    if audio_file and os.path.exists(audio_file):
        # Build FFmpeg command to combine image and audio with proper flags
        ffmpeg_cmd = [
            "ffmpeg",
            "-y",  # Overwrite output files without asking
            "-loop",
            "1",  # Loop the image for the duration of the audio
            "-i",
            image_file,  # Input image file
            "-i",
            audio_file,  # Input audio file
            "-c:v",
            "libx264",  # Use H.264 video codec
            "-tune",
            "stillimage",  # Tune for still image
            "-c:a",
            "aac",  # Use AAC audio codec
            "-b:a",
            "192k",  # Set audio bitrate
            "-pix_fmt",
            "yuv420p",  # Set pixel format for compatibility
            "-vf",
            f"scale={adjusted_width}:-2",  # Scale video width, height calculated to preserve aspect ratio
            "-shortest",  # Finish encoding when the shortest input ends
            slide_video_filename,  # Output video file
        ]
    else:
        logger.info(f"Creating silent video for slide {i+1}")
        # Build FFmpeg command to create a video from image for 5 seconds
        ffmpeg_cmd = [
            "ffmpeg",
            "-y",  # Overwrite output files without asking
            "-loop",
            "1",  # Loop the image
            "-i",
            image_file,  # Input image file
            "-c:v",
            "libx264",  # Use H.264 video codec
            "-t",
            "5",  # Set duration to 5 seconds
            "-pix_fmt",
            "yuv420p",  # Set pixel format for compatibility
            "-vf",
            f"scale={adjusted_width}:-2",  # Scale video width, height calculated to preserve aspect ratio
            slide_video_filename,  # Output video file
        ]

    try:
        subprocess.run(ffmpeg_cmd, check=True, capture_output=True)
        logger.info(f"Video for slide {i+1} created")
        return slide_video_filename
    except subprocess.CalledProcessError as e:
        logger.error(f"Error creating video for slide {i+1}: {e}")
        return None


class PPTXtoVideo:
    """A class to convert a PowerPoint presentation to a video using OpenAI TTS and FFmpeg."""

    def __init__(self, pptx_filename: str, force_flag: bool = False):
        """Initialize the PPTXtoVideo instance.

        :param pptx_filename: The path to the PowerPoint (.pptx) file.
        :param force_flag: If True, forces regeneration of all resources.
        """
        self.pptx_filename = pptx_filename
        self.force_flag = force_flag
        self.pptx_hash = self._compute_file_hash(pptx_filename)
        self.api_key = os.environ.get("OPENAI_API_KEY")
        if not self.api_key:
            raise ValueError("OPENAI_API_KEY not found in environment variables.")

        self.presentation = Presentation(pptx_filename)
        self.slides = self.presentation.slides
        self.slides_data = self._extract_slide_texts()
        self.video_files = []

        # Determine the output directory
        repo_root = self._get_repo_root()
        if repo_root:
            self.output_dir = os.path.join(repo_root, "pptx2video", "video-assets")
        else:
            # Fallback to using the PPTX file's directory
            self.output_dir = os.path.join(
                os.path.dirname(self.pptx_filename), "video-assets"
            )
        os.makedirs(self.output_dir, exist_ok=True)

        # Define filenames
        pptx_basename = os.path.splitext(os.path.basename(pptx_filename))[0]
        self.pdf_filename = os.path.join(self.output_dir, f"{pptx_basename}.pdf")
        self.output_file = os.path.join(
            os.path.dirname(pptx_filename), f"{pptx_basename}.mp4"
        )
        self.state_file = os.path.join(self.output_dir, "conversion_state.json")

        # Load or initialize state
        self.state = self._load_state()
        self.state_changed = False  # Flag to indicate whether state has changed

    def _get_repo_root(self) -> Optional[str]:
        """Get the root directory of the Git repository.

        :return: The path to the repository root or None if not in a Git repository.
        """
        try:
            result = subprocess.run(
                ["git", "rev-parse", "--show-toplevel"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                check=True,
            )
            repo_root = result.stdout.strip()
            logger.debug(f"Repository root found: {repo_root}")
            return repo_root
        except subprocess.CalledProcessError:
            logger.debug("Not in a Git repository or Git is not installed.")
            return None

    def _compute_file_hash(self, filename: str) -> str:
        """Compute a hash of the file contents.

        :param filename: Path to the file.
        :return: SHA256 hash of the file contents.
        """
        hasher = hashlib.sha256()
        with open(filename, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hasher.update(chunk)
        return hasher.hexdigest()

    def _extract_slide_texts(self) -> List[dict]:
        """Extract text from slide notes and compute hashes.

        :return: List of dicts with slide index, text, and text hash.
        """
        slides_data = []
        for i, slide in enumerate(self.slides):
            text = (
                slide.notes_slide.notes_text_frame.text.strip()
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text
                else ""
            )
            text_hash = hashlib.sha256(text.encode("utf-8")).hexdigest() if text else ""
            slides_data.append({"index": i, "text": text, "hash": text_hash})
            logger.debug(f"Slide {i+1} text hash: {text_hash} - text: {text[:50]}...")
        return slides_data

    def _load_state(self) -> dict:
        """Load the conversion state from a file.

        Handles missing file gracefully.
        :return: Dictionary containing the conversion state, or an empty dict if the file is missing.
        """
        initial_state = {
            "pptx_hash": "",
            "pdf_created": False,
            "slide_hashes": {},
            "slide_image_hashes": {},
            "slide_content_hashes": {},
            "final_video_hash": "",
        }
        try:
            with open(self.state_file, "r") as f:
                state = json.load(f)
                logger.debug("Conversion state loaded.")
        except (FileNotFoundError, json.JSONDecodeError):
            state = initial_state
            logger.debug(
                "No previous conversion state found or file corrupted. Starting fresh."
            )
        return state

    def _save_state(self):
        """Save the current conversion state to a file."""
        with open(self.state_file, "w") as f:
            json.dump(self.state, f, indent=4)
            logger.debug("Conversion state saved.")
        self.state_changed = False  # Reset the flag after saving

    def _convert_to_pdf(self):
        """Converts the .pptx file to a .pdf file using LibreOffice."""
        if (
            not self.force_flag
            and self.state.get("pdf_created", False)
            and os.path.exists(self.pdf_filename)
        ):
            logger.info(f"PDF is up to date, skipping conversion")
            return

        logger.info("Converting PPTX to PDF")
        os.makedirs(os.path.dirname(self.pdf_filename), exist_ok=True)

        cmd = [
            "libreoffice",
            "--headless",  # Run without GUI
            "--convert-to",
            "pdf",  # Convert to PDF format
            "--outdir",
            os.path.dirname(self.pdf_filename),  # Output directory
            self.pptx_filename,  # Input PPTX file
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        logger.debug(f"LibreOffice conversion output: {result.stdout}")
        if result.returncode != 0:
            logger.error(f"LibreOffice conversion errors: {result.stderr}")
            raise RuntimeError("Failed to convert PPTX to PDF")

        # Check if the PDF was created
        expected_pdf = os.path.join(
            os.path.dirname(self.pdf_filename),
            f"{os.path.splitext(os.path.basename(self.pptx_filename))[0]}.pdf",
        )
        if os.path.exists(expected_pdf):
            # Move the PDF to the desired location
            if expected_pdf != self.pdf_filename:
                os.replace(expected_pdf, self.pdf_filename)
            logger.info("PDF created successfully")
        else:
            logger.error("PDF file was not created")
            raise RuntimeError(f"Failed to create PDF file: {self.pdf_filename}")

        self.state["pdf_created"] = True
        self.state_changed = True

    def is_conversion_needed(self) -> Tuple[bool, dict]:
        """Check if conversion is needed and return details on what has changed.

        :return: Tuple of (conversion_needed: bool, changes: dict).
        """
        changes = {}
        conversion_needed = False

        # Check if PPTX file has changed
        if self.state.get("pptx_hash") != self.pptx_hash:
            changes["pptx_changed"] = True
            conversion_needed = True
        else:
            changes["pptx_changed"] = False

        # Verify that slides are up to date
        slides_changed = []
        for slide_data in self.slides_data:
            i = slide_data["index"]
            str_i = str(i)
            slide_notes_hash = slide_data["hash"]
            slide_image_hash = self.state["slide_image_hashes"].get(str_i, "")
            slide_content_hash = hashlib.sha256(
                (slide_notes_hash + slide_image_hash).encode()
            ).hexdigest()
            prev_slide_content_hash = self.state["slide_content_hashes"].get(str_i)
            if prev_slide_content_hash != slide_content_hash:
                slides_changed.append(i + 1)  # Slide numbers are 1-indexed
                conversion_needed = True

        changes["slides_changed"] = slides_changed

        # Compute current final video hash
        slide_indices = sorted(
            [int(i) for i in self.state.get("slide_content_hashes", {}).keys()]
        )
        concatenated_hash = "".join(
            [self.state["slide_content_hashes"].get(str(i), "") for i in slide_indices]
        )
        current_final_video_hash = hashlib.sha256(
            concatenated_hash.encode()
        ).hexdigest()

        if self.state.get("final_video_hash") != current_final_video_hash:
            changes["final_video_changed"] = True
            conversion_needed = True
        else:
            changes["final_video_changed"] = False

        return conversion_needed, changes

    async def create_videos(self):
        """Creates individual video files for each slide, combining slide images and TTS audio."""
        self._convert_to_pdf()

        # Convert PDF slides to images
        logger.info("Converting PDF slides to images")
        image_files = []
        images_to_generate = []
        image_hashes = {}

        # Collect slides that need image generation
        for i in range(len(self.slides)):
            image_file = os.path.join(self.output_dir, f"slide_{i}.png")
            image_files.append(image_file)
            str_i = str(i)
            if not os.path.exists(image_file) or self.force_flag:
                images_to_generate.append(
                    (self.pdf_filename, self.output_dir, i, self.force_flag)
                )
            else:
                # Compute image hash for existing image
                with open(image_file, "rb") as f:
                    image_hash = hashlib.sha256(f.read()).hexdigest()
                image_hashes[str_i] = image_hash

        if images_to_generate:
            with multiprocessing.Pool() as pool:
                # Convert slides to images in parallel
                results = pool.map(convert_slide_to_image, images_to_generate)
                for image_path, image_hash, i in results:
                    if image_hash:
                        str_i = str(i)
                        image_hashes[str_i] = image_hash
                    else:
                        logger.error(f"Failed to generate image for slide {i+1}")

        # Update state with current image hashes
        self.state["slide_image_hashes"] = image_hashes
        self.state_changed = True

        # Generate audio for slides
        logger.info("Generating audio for slides")
        audio_files = []
        tasks = []

        for slide_data in self.slides_data:
            i = slide_data["index"]
            text = slide_data["text"]
            text_hash = slide_data["hash"]
            audio_file = os.path.join(self.output_dir, f"voice_{i}.mp3")
            audio_files.append(audio_file)

            # Check if audio needs to be regenerated
            if (
                not self.force_flag
                and self.state["slide_hashes"].get(str(i)) == text_hash
                and os.path.exists(audio_file)
            ):
                logger.info(
                    f"Audio for slide {i + 1} is up to date, skipping generation"
                )
                continue

            # Limit the number of concurrent API calls
            task = asyncio.create_task(
                generate_audio_for_slide(
                    text,
                    text_hash,
                    self.output_dir,
                    i,
                    self.api_key,
                    self.state,
                    self.force_flag,
                )
            )
            tasks.append(task)

        # Await all audio generation tasks
        if tasks:
            await asyncio.gather(*tasks)
            self.state_changed = True

        # Save state after audio generation
        if self.state_changed:
            self._save_state()

        # Create videos for each slide
        logger.info("Creating videos for each slide")
        max_workers = min(multiprocessing.cpu_count(), len(self.slides))
        video_creation_args = []
        self.state.setdefault("slide_content_hashes", {})

        for i, (image_file, audio_file) in enumerate(zip(image_files, audio_files)):
            str_i = str(i)
            slide_notes_hash = self.slides_data[i]["hash"]
            slide_image_hash = self.state["slide_image_hashes"].get(str_i, "")
            slide_content_hash = hashlib.sha256(
                (slide_notes_hash + slide_image_hash).encode()
            ).hexdigest()
            prev_slide_content_hash = self.state["slide_content_hashes"].get(str_i)
            video_file = os.path.join(self.output_dir, f"video_{i}.mp4")
            if (
                not self.force_flag
                and prev_slide_content_hash == slide_content_hash
                and os.path.exists(video_file)
            ):
                logger.info(f"Video for slide {i + 1} is up to date, skipping")
                self.video_files.append(video_file)
                continue
            video_creation_args.append(
                (
                    image_file,
                    audio_file,
                    self.output_dir,
                    i,
                    slide_content_hash,
                    self.force_flag,
                )
            )

        if video_creation_args:
            with ProcessPoolExecutor(max_workers=max_workers) as executor:
                futures = {
                    executor.submit(create_video_for_slide, args): args
                    for args in video_creation_args
                }
                for future in futures:
                    args = futures[future]
                    i = args[3]
                    str_i = str(i)
                    slide_content_hash = args[4]
                    video_file = future.result()
                    if video_file:
                        self.video_files.append(video_file)
                        logger.info(f"Video for slide {i + 1} generated")
                        # Update state in the main process
                        self.state["slide_content_hashes"][str_i] = slide_content_hash
                        self.state_changed = True
                    else:
                        logger.error(f"Failed to create video for slide {i+1}")

        # Collect up-to-date video files
        for i in range(len(self.slides)):
            video_file = os.path.join(self.output_dir, f"video_{i}.mp4")
            if os.path.exists(video_file):
                if video_file not in self.video_files:
                    self.video_files.append(video_file)

        if self.state_changed:
            self._save_state()

    def combine_videos(self):
        """Concatenates individual slide videos into a final video."""
        # Compute final video hash
        slide_indices = sorted(
            [int(i) for i in self.state.get("slide_content_hashes", {}).keys()]
        )
        concatenated_hash = "".join(
            [self.state["slide_content_hashes"][str(i)] for i in slide_indices]
        )
        final_video_hash = hashlib.sha256(concatenated_hash.encode()).hexdigest()

        if (
            not self.force_flag
            and self.state.get("final_video_hash") == final_video_hash
            and os.path.exists(self.output_file)
        ):
            logger.info(f"Final video is up to date, skipping combination")
            return

        logger.info("Combining individual slide videos into final video")
        list_file = os.path.join(self.output_dir, "videos.txt")
        with open(list_file, "w") as f:
            for i in slide_indices:
                video_file = os.path.join(self.output_dir, f"video_{i}.mp4")
                if not os.path.exists(video_file):
                    raise FileNotFoundError(f"Missing video file: {video_file}")
                f.write(f"file '{video_file}'\n")

        # Build FFmpeg command to concatenate videos
        ffmpeg_cmd = [
            "ffmpeg",
            "-y",  # Overwrite output files without asking
            "-f",
            "concat",  # Indicate that we're providing a file list to concatenate
            "-safe",
            "0",  # Allow unsafe file paths
            "-i",
            list_file,  # Input list file
            "-c",
            "copy",  # Copy codec (no re-encoding)
            self.output_file,  # Output video file
        ]
        subprocess.run(ffmpeg_cmd, check=True, capture_output=True)
        logger.info(f"Final video created: {self.output_file}")

        # Update final video hash
        self.state["final_video_hash"] = final_video_hash
        self.state_changed = True
        if self.state_changed:
            self._save_state()

    async def convert(self):
        """Converts the PowerPoint presentation to a video."""
        logger.info(f"Starting conversion of {self.pptx_filename}")

        if self.force_flag:
            logger.info("Force flag detected, resetting state")
            self.state = {
                "pptx_hash": self.pptx_hash,
                "pdf_created": False,
                "slide_hashes": {},
                "slide_image_hashes": {},
                "slide_content_hashes": {},
                "final_video_hash": "",
            }
            self.state_changed = True
            self._save_state()
        else:
            # Update PPTX hash in state
            if self.state.get("pptx_hash") != self.pptx_hash:
                logger.info("PPTX file has changed, resetting state")
                self.state = {
                    "pptx_hash": self.pptx_hash,
                    "pdf_created": False,
                    "slide_hashes": {},
                    "slide_image_hashes": {},
                    "slide_content_hashes": {},
                    "final_video_hash": "",
                }
                self.state_changed = True
                self._save_state()

        conversion_needed, _ = self.is_conversion_needed()
        if not conversion_needed:
            logger.info("No conversion needed. Exiting.")
            return

        await self.create_videos()
        self.combine_videos()
        logger.info(f"Video created successfully: {self.output_file}")


async def main():
    """Main function to parse arguments and execute conversion."""
    parser = argparse.ArgumentParser(
        description="Convert a PowerPoint presentation to a video using OpenAI TTS and FFmpeg."
    )
    parser.add_argument(
        "pptx", type=str, help="The path to the PowerPoint (.pptx) file to convert."
    )
    parser.add_argument(
        "--force", action="store_true", help="Force regeneration of all resources."
    )
    parser.add_argument(
        "--changed",
        action="store_true",
        help="Output exactly what has changed that requires conversion.",
    )
    args = parser.parse_args()

    if args.changed:
        try:
            converter = PPTXtoVideo(args.pptx)
            conversion_needed, changes = converter.is_conversion_needed()
            if conversion_needed:
                print(json.dumps(changes, indent=4))
            else:
                print(json.dumps({"conversion_needed": False}, indent=4))
        except Exception as e:
            logger.exception(f"An error occurred while checking for changes: {e}")
            sys.exit(1)
        sys.exit(0)

    try:
        logger.info(f"Starting conversion process for: {args.pptx}")
        converter = PPTXtoVideo(args.pptx, force_flag=args.force)
        await converter.convert()
        logger.info("Conversion process completed successfully")
    except Exception as e:
        logger.exception(f"An unexpected error occurred: {e}")
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())
